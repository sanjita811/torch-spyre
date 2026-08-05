"""Generate the IBM-styled SSD presentation as an editable .pptx (python-pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- IBM / Carbon palette ---
INK   = RGBColor(0x16,0x16,0x16)
INK2  = RGBColor(0x52,0x51,0x4e)
MUTED = RGBColor(0x6f,0x6f,0x6f)
IBM   = RGBColor(0x0f,0x62,0xfe)
S2    = RGBColor(0xeb,0x68,0x34)   # orange
S3    = RGBColor(0x1b,0xaf,0x7a)   # aqua
GOOD  = RGBColor(0x19,0x80,0x38)
CRIT  = RGBColor(0xda,0x1e,0x28)
SURF  = RGBColor(0xf4,0xf4,0xf4)
LINE2 = RGBColor(0xc6,0xc6,0xc6)
WHITE = RGBColor(0xff,0xff,0xff)
FONT  = "IBM Plex Sans"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s

def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tb, tf

def setrun(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic
    f.name = FONT
    return r

def para(tf, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    return p

def kicker(s, text):
    tb, tf = box(s, 0.7, 0.55, 11, 0.5)
    p = para(tf, True); setrun(p, text, 13, IBM, bold=True)
    # accent underline bar
    bar = s.shapes.add_shape(1, Inches(0.72), Inches(0.5), Inches(0.28), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = IBM; bar.line.fill.background()

def title(s, text, y=1.0, size=34):
    tb, tf = box(s, 0.7, y, 11.9, 1.3)
    p = para(tf, True)
    setrun(p, text, size, INK, bold=True)
    return tb

def logo(s):
    # simple "IBM" wordmark bottom-left (8-bar logo is fiddly; wordmark reads clean)
    tb, tf = box(s, 0.7, 6.95, 2, 0.4)
    p = para(tf, True); setrun(p, "IBM", 15, INK, bold=True)

def footer(s, n):
    tb, tf = box(s, 11.6, 6.95, 1.4, 0.4)
    p = para(tf, True); p.alignment = PP_ALIGN.RIGHT
    setrun(p, f"{n} / 13", 11, MUTED)

def card(s, x, y, w, h, accent=LINE2):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = SURF; r.line.fill.background()
    r.shadow.inherit = False
    # top accent rule
    bar = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    bar.shadow.inherit = False
    return r

def bullets(tf, items, first=True):
    for k,(txt, col) in enumerate(items):
        p = para(tf, first and k==0)
        p.space_after = Pt(8)
        setrun(p, "▪  ", 15, col)
        # allow simple bold via **
        parts = txt.split("**")
        for j,seg in enumerate(parts):
            if seg: setrun(p, seg, 17, INK2, bold=(j%2==1))

def tile(s, x, y, w, h, value, label, vcolor=INK):
    card(s, x, y, w, h, accent=IBM)
    tb, tf = box(s, x+0.18, y+0.35, w-0.36, h-0.5)
    p = para(tf, True); setrun(p, value, 32, vcolor, bold=True)
    p2 = para(tf); p2.space_before = Pt(6)
    setrun(p2, label, 12, MUTED)

# ============ SLIDE 1 — TITLE ============
s = slide()
kicker(s, "TORCH-SPYRE · PROJECT REVIEW")
tb, tf = box(s, 0.7, 2.3, 11.5, 2)
p = para(tf, True); setrun(p, "Bringing State-Space Models", 46, INK, bold=True)
p = para(tf); setrun(p, "to the Spyre Accelerator", 46, INK, bold=True)
tb, tf = box(s, 0.7, 4.5, 10, 1)
p = para(tf, True)
setrun(p, "A production Mamba-2 SSD kernel, running fully on-device — the first "
          "structured-state-space sequence model on Spyre.", 20, MUTED)
tb, tf = box(s, 0.7, 5.7, 12, 0.5)
p = para(tf, True)
setrun(p, "Correctness validated vs. reference     ·     Optimized & profiled on-card"
          "     ·     Sequence length up to 32K", 14, INK2)
logo(s); footer(s, 1)

# ============ SLIDE 2 — WHY ============
s = slide(); kicker(s, "THE GOAL"); title(s, "Why state-space models matter")
tb, tf = box(s, 0.7, 2.0, 6.0, 4.5)
p = para(tf, True)
setrun(p, "Mamba-2 is a leading alternative to the Transformer. It processes long "
          "sequences in ", 19, INK2)
setrun(p, "linear time", 19, INK, bold=True)
setrun(p, " instead of quadratic — a structural advantage for long-context AI.", 19, INK2)
bullets(tf, [
    ("Transformer cost grows with the **square** of sequence length", IBM),
    ("SSMs stay **linear** — cheaper as context grows", S2),
    ("Used increasingly in frontier long-context & efficiency work", S3),
], first=False)
card(s, 7.1, 2.0, 5.5, 3.2, accent=IBM)
tb, tf = box(s, 7.35, 2.35, 5.0, 2.8)
p = para(tf, True); setrun(p, "The business case", 18, INK, bold=True)
p = para(tf); p.space_before=Pt(8)
setrun(p, "For Spyre to be a credible platform for modern AI, it must run the "
          "architectures customers want to deploy — not only Transformers.", 16, INK2)
p = para(tf); p.space_before=Pt(10)
setrun(p, "This project proves SSMs run on Spyre", 16, INK, bold=True)
setrun(p, ", and turns the accelerator's constraints into a concrete backend roadmap.", 16, INK2)
logo(s); footer(s, 2)

# ============ SLIDE 3 — CHALLENGE ============
s = slide(); kicker(s, "THE CHALLENGE"); title(s, "Spyre is not a GPU")
tb, tf = box(s, 0.7, 2.0, 6.0, 4.5)
p = para(tf, True)
setrun(p, "Spyre is a tiled, memory-oriented accelerator. Getting an SSM to run well "
          "means respecting hardware rules a GPU never imposes.", 19, INK2)
bullets(tf, [
    ("Data laid out in **128-byte blocks** — every reshape has a cost", IBM),
    ("No native **cumulative-sum**, no native attention op", S2),
    ("Per-core memory **limits** bound how big any tensor can be", S3),
    ("Half-precision throughout — **overflow** is a real constraint", IBM),
], first=False)
card(s, 7.1, 2.0, 5.5, 3.2)
tb, tf = box(s, 7.35, 2.35, 5.0, 2.8)
p = para(tf, True); setrun(p, "The reference is GPU-shaped", 18, INK, bold=True)
p = para(tf); p.space_before=Pt(8)
setrun(p, "The canonical Mamba-2 kernel is written for GPUs: giant fused operations, "
          "large chunk size, cumulative-sums everywhere.", 16, INK2)
p = para(tf); p.space_before=Pt(10)
setrun(p, "None of it maps directly.", 16, INK, bold=True)
setrun(p, " The work was re-deriving the same math as Spyre-legal operations — and "
          "finding where the hardware pushes back.", 16, INK2)
logo(s); footer(s, 3)

# ============ SLIDE 4 — APPROACH ============
s = slide(); kicker(s, "THE APPROACH"); title(s, "Decompose the algorithm into Spyre-legal matmuls")
tb, tf = box(s, 0.7, 1.9, 12, 1.1)
p = para(tf, True)
setrun(p, "The SSD algorithm splits a long sequence into chunks. Each stage — written "
          "as one big operation in the reference — becomes a chain of batched matrix "
          "multiplies with the decay folded into the operands.", 19, INK2)
stages = [("STAGE 1","Intra-chunk","local mixing within a chunk"),
          ("STAGE 2","Chunk states","summarize each chunk to a state"),
          ("STAGE 3","Inter-chunk scan","carry state across chunks"),
          ("STAGE 4","Combine","merge local + carried into output")]
x=0.7; w=2.85; gap=0.15
for i,(n,t,d) in enumerate(stages):
    card(s, x, 3.4, w, 1.7, accent=IBM)
    tb, tf = box(s, x+0.15, 3.55, w-0.3, 1.5)
    p=para(tf,True); setrun(p,n,11,IBM,bold=True)
    p=para(tf); p.space_before=Pt(3); setrun(p,t,17,INK,bold=True)
    p=para(tf); p.space_before=Pt(4); setrun(p,d,12,MUTED)
    if i<3:
        ar,af=box(s, x+w-0.02, 3.9, 0.35, 0.6)
        pp=para(af,True); pp.alignment=PP_ALIGN.CENTER; setrun(pp,"→",22,LINE2)
    x += w+gap
tb, tf = box(s, 0.7, 5.4, 12, 1.3)
p=para(tf,True)
setrun(p,"Key techniques: ",15,IBM,bold=True)
setrun(p,"cumulative-sum as a matrix multiply · decay folded into operands (no giant "
         "intermediate) · half-precision-safe exponentials · fused into one compiled kernel.",14,INK2)
p=para(tf); p.space_before=Pt(6)
setrun(p,"Design principle: ",15,S2,bold=True)
setrun(p,"the kernel is model-agnostic — it computes Y = SSM(X,A,B,C) only, so it is "
         "reusable across any SSM.",14,INK2)
logo(s); footer(s, 4)

# ============ SLIDE 5 — TESTING ============
s = slide(); kicker(s, "TESTING & VALIDATION"); title(s, "A CPU twin separates “our math” from “the backend”")
tb, tf = box(s, 0.7, 2.0, 6.2, 4.5)
bullets(tf, [
    ("**Reference oracle:** the verbatim Mamba-2 kernel on CPU as ground truth", IBM),
    ("**CPU mirror:** the same op-sequence as our kernel, in plain PyTorch — pinpoints "
     "whether a bug is in our formulation or the hardware", S2),
    ("**Device-faithful oracle:** a half-precision model that predicts Spyre behavior "
     "without a compile — so we tune correctness cheaply", S3),
], first=True)
p=para(tf); p.space_before=Pt(10)
setrun(p,"Acceptance metric: relative-L2 error < 0.05 — the realistic half-precision budget.",13,MUTED)
data=[("4e-4","fp32 formulation error — math matches reference"),
      ("9e-4","fp16 numerical floor — best in half precision"),
      ("0.0038","on-device error at default — near the floor"),
      ("< 0.05","every validated shape stays within budget")]
xs=[7.1,10.0]; ys=[2.0,4.05]
for i,(v,l) in enumerate(data):
    tile(s, xs[i%2], ys[i//2], 2.75, 1.85, v, l, vcolor=(GOOD if i==0 else IBM if i==1 else INK))
logo(s); footer(s, 5)

# ============ SLIDE 6 — ACCOMPLISHMENTS ============
s = slide(); kicker(s, "ACCOMPLISHMENTS"); title(s, "What works today")
acc=[("100%","On-device — full kernel on Spyre, no CPU fallback in the hot path",GOOD),
     ("32K","Max sequence length validated correct on-card",IBM),
     ("1","Fused kernel — intra + scan + combine in a single graph",INK),
     ("~2×","Kernel speedup from the factored-decay redesign",S2)]
x=0.7; w=2.95
for v,l,c in acc:
    tile(s, x, 2.1, w, 2.0, v, l, vcolor=c); x+=w+0.13
chips=["Fully fused single kernel","Accuracy at the fp16 floor","Auto-tuned per length",
       "Streaming initial-state","Safe fallback for large chunks","Reusable across SSMs"]
tb, tf = box(s, 0.7, 4.5, 12, 1.2)
p=para(tf,True)
for c in chips:
    setrun(p,f"  {c}  ",13,GOOD,bold=True); setrun(p,"   ",13,MUTED)
tb, tf = box(s, 0.7, 5.7, 12, 0.8)
p=para(tf,True)
setrun(p,"First structured-state-space model kernel to run on Spyre — from a GPU-only "
         "reference to a validated, tuned, on-device implementation.",15,INK2)
logo(s); footer(s, 6)

# ============ SLIDE 7 — PROFILING ============
s = slide(); kicker(s, "PROFILING"); title(s, "The kernel is memory-bound, not compute-bound")
tb, tf = box(s, 0.7, 2.0, 6.0, 4.5)
p=para(tf,True)
setrun(p,"Deep profiling gave the single most important insight of the project: the "
         "math is nearly free — ",19,INK2)
setrun(p,"moving data is the entire cost.",19,INK,bold=True)
p=para(tf); p.space_before=Pt(20)
setrun(p,"Ideal compute:  ",30,INK,bold=True); setrun(p,"~17 µs",30,IBM,bold=True)
p=para(tf)
setrun(p,"Runtime:  ",30,INK,bold=True); setrun(p,"milliseconds",30,S2,bold=True)
p=para(tf); p.space_before=Pt(16)
setrun(p,"Ideal matrix-multiply work is ~17 µs (the compiler's cycle model). The kernel "
         "runs orders of magnitude longer — the gap is entirely data movement.",13,MUTED)
card(s, 7.1, 2.0, 5.5, 3.6, accent=IBM)
tb, tf = box(s, 7.35, 2.35, 5.0, 3.2)
p=para(tf,True); setrun(p,"Why this matters strategically",18,INK,bold=True)
p=para(tf); p.space_before=Pt(8)
setrun(p,"It redirected the entire optimization effort. Chasing compute utilization "
         "would have been wasted work.",16,INK2)
p=para(tf); p.space_before=Pt(8)
setrun(p,"Every real win came from cutting data movement",16,INK,bold=True)
setrun(p," — fewer layout conversions, smaller intermediates, sharing tensors across ops.",16,INK2)
logo(s); footer(s, 7)

# ============ SLIDE 8 — OPTIMIZATIONS ============
s = slide(); kicker(s, "OPTIMIZATIONS"); title(s, "Every win cut data movement")
tb, tf = box(s, 0.7, 2.0, 6.2, 4.5)
bullets(tf, [
    ("**Factored decay:** fold decay into operands, not a giant matrix — killed a 34 MB "
     "intermediate, **~2× faster**", IBM),
    ("**Shared layouts:** reuse one prepared tensor across two matmuls — measured to the "
     "minimum number of conversions", S2),
    ("**Asymmetric mask build:** broadcast one operand in-kernel — **halved** a 268 MB "
     "transfer to 134 MB", S3),
    ("**Tiling sweep:** measured optimal core-tiling on-card — **~34% faster** than the "
     "previous default", IBM),
], first=True)
card(s, 7.1, 2.0, 5.5, 3.0)
tb, tf = box(s, 7.35, 2.35, 5.0, 2.6)
p=para(tf,True); setrun(p,"Rigor over intuition",18,INK,bold=True)
p=para(tf); p.space_before=Pt(8)
setrun(p,"On this hardware, counting operations does not predict speed",16,INK,bold=True)
setrun(p," — a cheaper-looking layout can be slower.",16,INK2)
p=para(tf); p.space_before=Pt(8)
setrun(p,"So every change was measured on-device, and several plausible ideas were "
         "rejected when the hardware disagreed.",16,INK2)
logo(s); footer(s, 8)

# ============ SLIDE 9 — CHUNK SWEEP ============
s = slide(); kicker(s, "OPTIMIZATION · MEASURED"); title(s, "Finding: GPU's chunk size is wrong for Spyre")
tb, tf = box(s, 0.7, 2.0, 5.8, 1.5)
p=para(tf,True)
setrun(p,"GPUs use a large chunk size (256). On memory-bound Spyre, the ",19,INK2)
setrun(p,"smallest",19,INK,bold=True)
setrun(p," chunk is fastest — the opposite conclusion.",19,INK2)
# simple bar chart as shapes
chart=[("chunk 64 (best)",928,IBM),("chunk 128",1777,S2),("chunk 256 (GPU)",4402,RGBColor(0xed,0xa1,0x00))]
bx, by, bw, maxh = 0.9, 6.2, 1.4, 3.3; maxv=4402
for i,(lab,val,col) in enumerate(chart):
    h=maxh*val/maxv
    r=s.shapes.add_shape(1, Inches(bx+i*1.9), Inches(by-h), Inches(bw), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb=col; r.line.fill.background(); r.shadow.inherit=False
    tb,tf=box(s, bx+i*1.9-0.2, by-h-0.4, bw+0.4, 0.35)
    p=para(tf,True); p.alignment=PP_ALIGN.CENTER; setrun(p,f"{val} ms",13,INK,bold=True)
    tb,tf=box(s, bx+i*1.9-0.25, by+0.05, bw+0.5, 0.5)
    p=para(tf,True); p.alignment=PP_ALIGN.CENTER; setrun(p,lab,12,INK2)
tb,tf=box(s, 0.7, 6.75, 6, 0.4); p=para(tf,True)
setrun(p,"Sequence length 4,096. Lower is faster. Measured on-card.",12,MUTED)
card(s, 7.1, 2.0, 5.5, 3.0)
tb, tf = box(s, 7.35, 2.35, 5.0, 2.6)
p=para(tf,True); setrun(p,"Why smaller wins",18,INK,bold=True)
p=para(tf); p.space_before=Pt(8)
setrun(p,"Within-chunk work grows with the square of chunk size. Because the kernel is "
         "memory-bound, that intermediate dominates — the smallest legal chunk moves the "
         "least data.",16,INK2)
p=para(tf); p.space_before=Pt(8)
setrun(p,"Takeaway: ",16,INK,bold=True)
setrun(p,"hardware-specific tuning beats porting GPU defaults.",16,INK2)
logo(s); footer(s, 9)

# ============ SLIDE 10 — ENVELOPE ============
s = slide(); kicker(s, "LONG-CONTEXT ENVELOPE"); title(s, "Correct & running up to 32K tokens")
env=[("4K",0.0039,1),("8K",0.0047,1),("16K",0.0057,1),("32K",0.0109,1),("64K",None,0)]
bx,by,bw,maxh=0.9,5.6,0.95,3.2; maxv=0.05
for i,(lab,val,ok) in enumerate(env):
    xx=bx+i*1.15
    if ok:
        h=maxh*val/maxv
        r=s.shapes.add_shape(1, Inches(xx), Inches(by-h), Inches(bw), Inches(h))
        r.fill.solid(); r.fill.fore_color.rgb=GOOD; r.line.fill.background(); r.shadow.inherit=False
        tb,tf=box(s, xx-0.25, by-h-0.35, bw+0.5, 0.3)
        p=para(tf,True); p.alignment=PP_ALIGN.CENTER; setrun(p,f"{val:.4f}",11,INK,bold=True)
    else:
        h=1.0
        r=s.shapes.add_shape(1, Inches(xx), Inches(by-h), Inches(bw), Inches(h))
        r.fill.background(); r.line.color.rgb=CRIT; r.line.width=Pt(1.5); r.shadow.inherit=False
        tb,tf=box(s, xx-0.25, by-h+0.25, bw+0.5, 0.4)
        p=para(tf,True); p.alignment=PP_ALIGN.CENTER; setrun(p,"✗",20,CRIT,bold=True)
    tb,tf=box(s, xx-0.25, by+0.05, bw+0.5, 0.35)
    p=para(tf,True); p.alignment=PP_ALIGN.CENTER; setrun(p,lab,13,(CRIT if not ok else INK2))
tb,tf=box(s,0.7,6.05,6,0.4); p=para(tf,True)
setrun(p,"Green = validated on-device.  Budget = 0.05.  Bars scaled to budget.",12,MUTED)
tb, tf = box(s, 7.1, 2.0, 5.5, 4.3)
p=para(tf,True)
setrun(p,"The maximum sequence length ",19,INK2)
setrun(p,"doubled this cycle",19,INK,bold=True)
setrun(p," — from 16K to 32K — by re-measuring a hardware limit against the latest backend.",19,INK2)
bullets(tf,[
    ("Each length uses its own measured-best chunk size",IBM),
    ("Error stays within budget across the whole range",S2),
    ("**64K+** is blocked by a specific, identified backend limitation — not by our math",S3),
], first=False)
logo(s); footer(s, 10)

# ============ SLIDE 11 — FUTURE ============
s = slide(); kicker(s, "FUTURE WORK"); title(s, "One backend fix unlocks the next frontier")
tb, tf = box(s, 0.7, 2.0, 6.0, 4.5)
p=para(tf,True)
setrun(p,"Beyond 32K, one tensor grows past the per-core memory limit. The algorithmic "
         "fix is known and CPU-validated — it needs one backend capability.",19,INK2)
bullets(tf,[
    ("**Sub-quadratic scan** — validated correct on CPU; the math is ready",IBM),
    ("**The blocker:** reshaping a tensor into blocks corrupts its layout on-device",S2),
    ("**The ask:** correct layout across a splitting reshape — one well-scoped capability",S3),
], first=False)
card(s, 7.1, 2.0, 5.5, 3.4, accent=IBM)
tb, tf = box(s, 7.35, 2.35, 5.0, 3.0)
p=para(tf,True); setrun(p,"Turning a kernel into a roadmap",18,INK,bold=True)
p=para(tf); p.space_before=Pt(8)
setrun(p,"This project produced a prioritized list of backend capabilities, each with a "
         "minimal reproducer and a clear payoff:",16,INK2)
for txt,col in [("Splitting-reshape layout → unlocks 64K+",CRIT),
                ("On-chip residency → cuts the memory tax",S2),
                ("Layout-sharing across ops → general speedup",S2)]:
    p=para(tf); p.space_before=Pt(6); setrun(p,"▪  ",14,col); setrun(p,txt,14,INK2)
logo(s); footer(s, 11)

# ============ SLIDE 12 — IMPACT ============
s = slide(); kicker(s, "IMPACT & LONG-TERM GOAL"); title(s, "From one kernel to a platform capability")
card(s, 0.7, 2.0, 5.9, 2.3, accent=IBM)
tb, tf = box(s, 0.95, 2.3, 5.4, 1.9)
p=para(tf,True); setrun(p,"Impact today",17,INK,bold=True)
for txt,col in [("SSMs are **proven viable** on Spyre",IBM),
                ("A validated, tuned, reusable kernel others build on",S2),
                ("A concrete, evidence-backed backend roadmap",S3)]:
    p=para(tf); p.space_before=Pt(5); setrun(p,"▪  ",13,col)
    for j,seg in enumerate(txt.split("**")):
        if seg: setrun(p,seg,15,INK2,bold=(j%2==1))
card(s, 6.8, 2.0, 5.8, 2.3)
tb, tf = box(s, 7.05, 2.3, 5.3, 1.9)
p=para(tf,True); setrun(p,"Long-term goal",17,INK,bold=True)
for txt,col in [("**Full Mamba-2 layer** — projections, conv, norm, decode",IBM),
                ("End-to-end **SSM inference** for long-context serving",S2),
                ("Backend fixes that speed up the whole platform",S3)]:
    p=para(tf); p.space_before=Pt(5); setrun(p,"▪  ",13,col)
    for j,seg in enumerate(txt.split("**")):
        if seg: setrun(p,seg,15,INK2,bold=(j%2==1))
roadmap=[("DONE","SSD core kernel","on-device, tuned, 32K"),
         ("NEXT","Full layer","proj · conv · norm · gate"),
         ("THEN","Decode step","recurrent inference"),
         ("GOAL","SSM serving","long-context, on Spyre")]
x=0.7; w=2.95
for i,(n,t,d) in enumerate(roadmap):
    card(s, x, 4.8, w, 1.5, accent=IBM)
    tb,tf=box(s, x+0.15, 4.95, w-0.3, 1.3)
    p=para(tf,True); setrun(p,n,11,IBM,bold=True)
    p=para(tf); p.space_before=Pt(3); setrun(p,t,16,INK,bold=True)
    p=para(tf); p.space_before=Pt(3); setrun(p,d,11,MUTED)
    if i<3:
        ar,af=box(s, x+w-0.05, 5.2, 0.35, 0.5)
        pp=para(af,True); pp.alignment=PP_ALIGN.CENTER; setrun(pp,"→",20,LINE2)
    x+=w+0.13
logo(s); footer(s, 12)

# ============ SLIDE 13 — CLOSE ============
s = slide(); kicker(s, "SUMMARY")
tb, tf = box(s, 0.7, 2.1, 12, 2)
p=para(tf,True); setrun(p,"SSMs run on Spyre —",44,INK,bold=True)
p=para(tf); setrun(p,"correct, tuned, up to 32K.",44,INK,bold=True)
tb, tf = box(s, 0.7, 4.4, 5.9, 2)
p=para(tf,True)
setrun(p,"We took a GPU-only architecture, re-derived it for Spyre's constraints, "
         "validated it to the precision floor, and optimized it to its memory-bound ceiling.",19,INK2)
tb, tf = box(s, 6.8, 4.4, 5.8, 2)
p=para(tf,True)
setrun(p,"The path to full SSM serving is clear — and the backend work it needs is "
         "identified, reproduced, and prioritized.",19,INK2)
tb, tf = box(s, 0.7, 6.2, 8, 0.5)
p=para(tf,True); setrun(p,"Thank you   ·   Questions welcome",15,IBM,bold=True)
logo(s); footer(s, 13)

prs.save("ssd_presentation.pptx")
print("saved ssd_presentation.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
